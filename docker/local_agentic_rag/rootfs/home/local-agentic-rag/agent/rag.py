import os
import sys
import shutil
import zipfile
import xml.etree.ElementTree as ET
import lancedb
import pandas as pd
from sentence_transformers import SentenceTransformer

# Disable tokenizer parallelism warning commonly seen in CLI loops
os.environ["TOKENIZERS_PARALLELISM"] = "false"

# Directory paths (Executes from the repository root)
INPUT_DIR = "/app/_raw_input"
PROCESSED_DIR = os.path.join(INPUT_DIR, ".processed")

# Supported formats list
SUPPORTED_EXTENSIONS = (
    '.txt', '.md', '.pdf', '.docx', '.pptx',
    '.csv', '.xlsx', '.xls', '.ods',
    '.html', '.htm', '.odt', '.odp'
)

# 1. Get the query or command from CMD
if len(sys.argv) < 2:
    sys.exit(0)

command_action = sys.argv[1].lower()
query = " ".join(sys.argv[1:])

# 2. Initialize LanceDB database connection
db = lancedb.connect("/home/local-agentic-rag/agent/.lancedb")

# ACTION: WIPE
if command_action == "wipe":
    # Robust wipe using try-except instead of list_tables()
    try:
        db.drop_table("documents")
    except Exception:
        pass

    if os.path.exists(PROCESSED_DIR):
        shutil.rmtree(PROCESSED_DIR)
    os.makedirs(PROCESSED_DIR, exist_ok=True)

    print("Database and processed files wiped successfully.")
    sys.exit(0)

# Initialize local model (Lazy load: skipped if action was 'wipe')
try:
    # Try loading from local cache first (no network requests, no HF warnings, much faster)
    model = SentenceTransformer('all-MiniLM-L6-v2', local_files_only=True)
except Exception:
    # If not found locally, download it (will show the warning only once)
    model = SentenceTransformer('all-MiniLM-L6-v2')

# Simple function definition for text chunking
def chunk_text(text, chunk_size=500, overlap=100):
    words = text.split()
    chunks = []
    for i in range(0, len(words), chunk_size - overlap):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
    return chunks

# Function for extracting text from multiple file formats
def extract_text_from_file(file_path):
    ext = file_path.lower().split('.')[-1]

    try:
        # Text and Markdown
        if ext in ['txt', 'md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        # PDF
        elif ext == 'pdf':
            import pypdf
            reader = pypdf.PdfReader(file_path)
            return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

        # MS Word
        elif ext == 'docx':
            import docx
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        # MS PowerPoint
        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)

        # CSV
        elif ext == 'csv':
            df = pd.read_csv(file_path)
            return df.to_string(index=False)

        # Excel and LibreOffice Calc
        elif ext in ['xlsx', 'xls', 'ods']:
            engine = 'odf' if ext == 'ods' else 'openpyxl'
            df = pd.read_excel(file_path, engine=engine)
            return df.to_string(index=False)

        # HTML
        elif ext in ['html', 'htm']:
            from bs4 import BeautifulSoup
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                return soup.get_text(separator='\n', strip=True)

        # LibreOffice Writer & Impress (lossless workaround via ZIP/XML format)
        elif ext in ['odt', 'odp']:
            with zipfile.ZipFile(file_path, 'r') as z:
                content_xml = z.read('content.xml')
                tree = ET.fromstring(content_xml)
                text = []
                for elem in tree.iter():
                    if elem.text and elem.text.strip():
                        text.append(elem.text.strip())
                return "\n".join(text)

        return ""
    except Exception as e:
        # Silently ignore reading errors in CLI background (e.g., corrupted file, empty rows)
        return ""

# 3. INGESTION OF NEW FILES
new_chunks = []

# Scan files in the main _raw_input directory
if os.path.exists(INPUT_DIR):
    # Create .processed directory if it doesn't exist
    os.makedirs(PROCESSED_DIR, exist_ok=True)

    for filename in os.listdir(INPUT_DIR):
        file_path = os.path.join(INPUT_DIR, filename)

        # Ignore subfolders and process only supported extensions
        if os.path.isfile(file_path) and filename.lower().endswith(SUPPORTED_EXTENSIONS):
            try:
                content = extract_text_from_file(file_path)

                if content.strip():
                    file_chunks = chunk_text(content)
                    new_chunks.extend(file_chunks)

                # Move file to archive (even if empty, to avoid loops)
                shutil.move(file_path, os.path.join(PROCESSED_DIR, filename))
            except Exception as e:
                pass

# If new documents are found, generate vectors and add them to the database
if new_chunks:
    data = []
    for chunk in new_chunks:
        vector = model.encode(chunk).tolist()
        data.append({"vector": vector, "text": chunk})

    # Robust save to table using try-except instead of checking lists
    try:
        # Attempt to open and append first
        tbl = db.open_table("documents")
        tbl.add(data)
    except Exception:
        # If opening fails, table likely doesn't exist yet - let's create it
        try:
            tbl = db.create_table("documents", data=data)
        except ValueError:
            # Ultimate fallback in case of race conditions
            tbl = db.open_table("documents")
            tbl.add(data)
else:
    # If no new files, simply open the existing database
    try:
        tbl = db.open_table("documents")
    except Exception:
        tbl = None

# ACTION: SYNC
if command_action == "sync":
    print("Synchronization complete. New files processed.")
    sys.exit(0)

# 4. RETRIEVAL
if tbl is not None:
    query_vector = model.encode(query).tolist()
    # Get top 3 matching fragments for better context
    results = tbl.search(query_vector).limit(3).to_pandas()

    if not results.empty:
        # Combine results into one text block separated by newlines
        context = "\n---\n".join(results['text'].tolist())
        print(context)
    else:
        print("No matching documents found in the database.")
else:
    print("Document database is empty. Drop files into the ./_raw_input folder.")